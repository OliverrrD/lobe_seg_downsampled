import os
from pptx import Presentation
from pptx.util import Inches, Pt

def save_pptx(png_file_list, pptx_path):
    prs = Presentation()
    layout_image = prs.slide_layouts[6]
    left = top = Inches(0)
    height = Inches(5.5)
    font_size = Pt(10)
    left_text = Inches(1)
    top_text = Inches(6.6)

    for png_file_path in png_file_list:
        print(f'Add {png_file_path}')
        slide = prs.slides.add_slide(layout_image)
        slide.shapes.add_picture(png_file_path, left, top, height=height)

        txBox = slide.shapes.add_textbox(left_text, top_text, Inches(1), Inches(1))
        tf = txBox.text_frame
        tf.text = os.path.basename(png_file_path.replace('.png', ''))
        # tf.font.size = font_size

    print(f'Save to {pptx_path}')
    prs.save(pptx_path)